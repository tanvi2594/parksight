"""ParkSight pitch deck — detailed, professional, with REAL embedded screenshots + speaker notes."""
import json
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent
M = json.load(open(ROOT / "outputs" / "metrics.json"))
f, ho, v = M["forecast"], M["holdout"], M["validation"]
cost, o, roi = M["cost"], M["offenders"], M["roi"]
cov = M.get("coverage", {}); emg = M.get("emerging", {})

FONT = "Segoe UI"
NAVY = RGBColor(0x0F, 0x1D, 0x35); NAVY2 = RGBColor(0x24, 0x3A, 0x60)
BLUE = RGBColor(0x25, 0x63, 0xEB); SKY = RGBColor(0x9D, 0xC2, 0xFF)
INK = RGBColor(0x1F, 0x29, 0x37); MUTE = RGBColor(0x6B, 0x7A, 0x90)
WHITE = RGBColor(0xFF, 0xFF, 0xFF); GREEN = RGBColor(0x05, 0x96, 0x69)
RED = RGBColor(0xDC, 0x26, 0x26); VIO = RGBColor(0x7C, 0x3A, 0xED); AMBER = RGBColor(0xD9, 0x77, 0x06)
LIGHT = RGBColor(0xF4, 0xF7, 0xFB); LINE = RGBColor(0xE3, 0xE8, 0xF0)
CARD = RGBColor(0xFB, 0xFC, 0xFE); MINT = RGBColor(0xEC, 0xFD, 0xF5); MINTL = RGBColor(0xA7, 0xF3, 0xD0)

# ── screenshot filenames ──────────────────────────────────────────────────
SHOT = lambda t: f"Screenshot 2026-06-21 {t}.png"
HERO = SHOT("152248"); CMAP = SHOT("152310"); SAT = SHOT("152347"); CITY3D = SHOT("152541")
POI = SHOT("152554"); LIVE = SHOT("152619"); PRIO = SHOT("152632"); PROUTE = SHOT("152641")
ROLLUP = SHOT("152652"); DRILL = SHOT("152658"); FCAST = SHOT("152710"); FWIN = SHOT("152719")
OFF = SHOT("152728"); TREND = SHOT("152735"); TLAPSE = SHOT("152743"); GAP = SHOT("152803")
BIAS = SHOT("152810"); OSM = SHOT("152820"); VAL = SHOT("152831"); VALC = SHOT("152838")
METH = SHOT("152847")

CLEAN = ROOT / "snapshots" / "_deck"; CLEAN.mkdir(exist_ok=True)
_cache = {}


def clean(name, top=0.10, bottom=0.052, left=0.0, right=0.0):
    """Crop browser chrome (top) + Windows taskbar (bottom), cache, return path + (w,h)."""
    key = (name, top, bottom, left, right)
    if key in _cache:
        return _cache[key]
    im = Image.open(ROOT / "snapshots" / name)
    w, h = im.size
    crop = im.crop((int(w * left), int(h * top), int(w * (1 - right)), int(h * (1 - bottom))))
    dst = CLEAN / (name.replace(" ", "_").replace(".png", f"_{int(left*100)}{int(right*100)}.png"))
    crop.save(dst)
    _cache[key] = (str(dst), crop.size)
    return _cache[key]


prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]; SW, SH = prs.slide_width, prs.slide_height
_page = [0]


def _bg(s, c):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH); r.fill.solid()
    r.fill.fore_color.rgb = c; r.line.fill.background(); r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)


def box(s, x, y, w, h):
    tf = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.word_wrap = True; return tf


def txt(tf, t, size=16, color=INK, bold=False, align=PP_ALIGN.LEFT, space=5, mono=False, first=False):
    p = tf.paragraphs[0] if (first or (len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs)) else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space)
    r = p.add_run(); r.text = str(t)
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    r.font.name = "Consolas" if mono else FONT
    return p


def rect(s, x, y, w, h, c, shape=MSO_SHAPE.RECTANGLE, line=None, lw=1.0):
    sh = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.shadow.inherit = False
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
    return sh


def note(s, t):
    s.notes_slide.notes_text_frame.text = t


def footer(s):
    _page[0] += 1
    rect(s, 0, 7.18, 13.333, 0.012, LINE)
    txt(box(s, 0.6, 7.2, 9, 0.3), "ParkSight  ·  Gridlock 2.0  ·  Theme 1 — Parking-Induced Congestion", 9, MUTE)
    txt(box(s, 12.2, 7.2, 0.7, 0.3), str(_page[0]), 9, MUTE, align=PP_ALIGN.RIGHT)


def content(kicker, title, sub=None):
    s = prs.slides.add_slide(BLANK); _bg(s, WHITE)
    rect(s, 0.6, 0.52, 0.16, 0.62, BLUE)
    txt(box(s, 0.95, 0.48, 11.6, 0.34), kicker, 12, BLUE, True)
    txt(box(s, 0.95, 0.82, 11.7, 0.66), title, 25, INK, True)
    rect(s, 0.95, 1.5, 11.45, 0.02, LINE)
    if sub:
        txt(box(s, 0.95, 1.56, 11.6, 0.4), sub, 12.5, MUTE)
    footer(s)
    return s


def divider(num, kicker, title):
    s = prs.slides.add_slide(BLANK); _bg(s, NAVY)
    txt(box(s, 0.9, 1.7, 5, 1.8), num, 96, NAVY2, True)
    rect(s, 0.98, 3.95, 2.2, 0.06, BLUE)
    txt(box(s, 0.98, 4.15, 11, 0.4), kicker, 13, SKY, True)
    txt(box(s, 0.98, 4.55, 11.5, 1.0), title, 30, WHITE, True)
    return s


def card(s, x, y, w, val, label, color=BLUE, h=1.55):
    rect(s, x, y, w, h, LIGHT, MSO_SHAPE.ROUNDED_RECTANGLE, line=LINE, lw=1)
    rect(s, x + 0.12, y + 0.12, w - 0.24, 0.08, color)
    tf = box(s, x + 0.18, y + 0.28, w - 0.32, h - 0.36)
    txt(tf, val, 21, color, True, first=True); txt(tf, label, 10.5, MUTE, space=0)


def bullets(s, x, y, w, h, items, size=15, space=8):
    tf = box(s, x, y, w, h)
    for i, t in enumerate(items):
        if t == "":
            txt(tf, "", 6, MUTE, space=2, first=(i == 0)); continue
        txt(tf, t, size, INK if not t.startswith("›") else MUTE, bold=t.endswith(":"), space=space, first=(i == 0))
    return tf


def img(s, x, y, w, h, name, caption=None, accent=BLUE, pad=0.14, **crop):
    """Place a cleaned screenshot fit-contained on a bordered card, optional caption."""
    path, (iw, ih) = clean(name, **crop)
    rect(s, x, y, w, h, CARD, MSO_SHAPE.ROUNDED_RECTANGLE, line=LINE, lw=1)
    rect(s, x + 0.14, y + 0.13, 0.5, 0.08, accent)
    cap_h = 0.32 if caption else 0.0
    fw, fh = w - 2 * pad, h - 2 * pad - cap_h
    ar = iw / ih
    if fw / fh > ar:
        dh = fh; dw = dh * ar
    else:
        dw = fw; dh = dw / ar
    ix = x + (w - dw) / 2
    iy = y + pad + (fh - dh) / 2
    pic = s.shapes.add_picture(path, Inches(ix), Inches(iy), Inches(dw), Inches(dh))
    pic.line.color.rgb = LINE; pic.line.width = Pt(0.75); pic.shadow.inherit = False
    if caption:
        txt(box(s, x + 0.1, y + h - cap_h - 0.04, w - 0.2, cap_h), caption, 9.5, MUTE, align=PP_ALIGN.CENTER, space=0)


def table(s, x, y, w, data, col_w, row_h=0.34, head_h=0.4, fs=11.5):
    rows, cols = len(data), len(data[0]); h = head_h + row_h * (rows - 1)
    gt = s.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    gt.first_row = False; gt.horz_banding = False
    for ci in range(cols):
        gt.columns[ci].width = Inches(col_w[ci])
    for ri, row in enumerate(data):
        gt.rows[ri].height = Inches(head_h if ri == 0 else row_h)
        for ci, val in enumerate(row):
            cell = gt.cell(ri, ci); cell.margin_left = Inches(0.1); cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03); cell.vertical_anchor = MSO_ANCHOR.MIDDLE; cell.text = str(val)
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            rn = p.runs[0]; rn.font.name = FONT; rn.font.size = Pt(fs)
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY; rn.font.color.rgb = WHITE; rn.font.bold = True
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT; rn.font.color.rgb = INK
                if ci == 1:
                    rn.font.bold = True; rn.font.color.rgb = BLUE
    return gt


# ════════════════════════════ COVER
s = prs.slides.add_slide(BLANK); _bg(s, NAVY)
rect(s, 0, 0, 0.22, 7.5, BLUE); rect(s, 0.9, 3.18, 3.4, 0.06, BLUE)
txt(box(s, 0.85, 1.9, 11.5, 1.1), "ParkSight", 56, WHITE, True)
txt(box(s, 0.9, 3.35, 11.6, 0.7), "Parking-Induced Congestion Intelligence", 23, SKY)
txt(box(s, 0.9, 4.25, 11.4, 0.9),
    "An AI decision-support system that detects illegal-parking hotspots, quantifies their impact on "
    "traffic flow, forecasts risk, and produces an optimised, road-routed enforcement plan.", 14,
    RGBColor(0xC7, 0xD4, 0xEA))
txt(box(s, 0.9, 6.35, 11.6, 0.7),
    "Gridlock 2.0  ·  Theme 1  ·  Built on Bengaluru Traffic Police (ASTraM) data + MapmyIndia (Mappls) maps\n"
    "Dataset: jan to may police violation_anonymized791b166.csv   ·   validated vs astram_event_data.csv",
    11, RGBColor(0x8A, 0xA0, 0xC0))
note(s, "Opening line: 'Illegal parking is one of the most fixable causes of urban congestion — yet "
        "enforcement today is blind and reactive. ParkSight turns the traffic police's own ~300,000 violation "
        "records into a live, validated, deployable plan that says exactly where, when, and how to enforce.' "
        "Stress: built on the partners' real data and MapmyIndia maps — deployable, not a toy.")

# ════════════════════════════ 01
divider("01", "PROBLEM & APPROACH", "Enforcement is blind, reactive, un-prioritised")
note(prs.slides[-1], "Section 1 sets up the operational pain and our framing: turn logs into a decision loop.")

s = content("THE PROBLEM", "Four gaps in parking enforcement today",
            "On-street & spillover illegal parking near markets, metros and events chokes carriageways.")
for (title, body, col), x in zip(
        [("Reactive", "Patrol-based; officers respond to complaints, not data.", RED),
         ("Invisible", "No heatmap of where violations actually hurt traffic flow.", AMBER),
         ("Un-prioritised", "No objective way to allocate limited staff city-wide.", VIO),
         ("Un-measured", "No feedback loop to learn whether an action worked.", BLUE)],
        [0.95, 3.95, 6.95, 9.95]):
    rect(s, x, 2.2, 2.85, 2.7, LIGHT, MSO_SHAPE.ROUNDED_RECTANGLE, line=LINE, lw=1)
    rect(s, x + 0.18, 2.42, 0.5, 0.1, col)
    tf = box(s, x + 0.18, 2.62, 2.55, 2.2)
    txt(tf, title, 16, INK, True, first=True); txt(tf, body, 12, MUTE, space=0)
txt(box(s, 0.95, 5.3, 11.6, 0.7),
    "Goal: convert raw violation logs into a ranked, validated, deployable enforcement plan.", 15, INK, True)
note(s, "Walk the four cards left-to-right. Land the goal line: we're not just visualising data — we're "
        "producing an ACTION plan a DCP can deploy tomorrow. These four gaps map to our five-stage loop.")

s = content("THE SOLUTION", "A closed, end-to-end decision loop")
x = 0.95
for i, t in enumerate(["Detect\nhotspots", "Quantify\nimpact", "Forecast\nrisk", "Optimise\n& deploy", "Validate\n& learn"]):
    c = rect(s, x, 2.7, 2.55, 1.4, BLUE if i % 2 == 0 else NAVY, MSO_SHAPE.CHEVRON)
    c.text_frame.word_wrap = True; txt(c.text_frame, t, 14, WHITE, True, PP_ALIGN.CENTER, first=True); x += 2.3
bullets(s, 0.95, 4.7, 11.7, 1.8, [
    "Each stage is a working module in a live Streamlit command center (9 interactive views).",
    "The loop re-runs monthly, so any enforcement action can be measured — the post-event learning the brief asks for.",
], size=15, space=9)
note(s, "Emphasise CLOSED loop — most teams stop at 'detect/visualise'. We forecast, optimise deployment, "
        "validate against independent data, and re-measure monthly. That last step is exactly the 'post-event "
        "learning system' the problem statement asks for.")

s = content("ARCHITECTURE", "Reproducible pipeline → artifacts → live app")
rect(s, 0.95, 1.8, 11.45, 3.55, LIGHT, MSO_SHAPE.ROUNDED_RECTANGLE, line=LINE, lw=1)
mono = box(s, 1.2, 2.0, 11.0, 3.2)
for i, t in enumerate([
    "jan to may police violation_anonymized791b166.csv (298K) ─► pipeline.py   clean · severity · H3",
    "                             ├─► hotspots.py   H3 cells + Congestion-Impact Score ─► DBSCAN zones",
    "                             ├─► forecast.py   HistGB + LightGBM ensemble · hold-out · intervals",
    "                             ├─► analytics.py  offenders · emerging · coverage-gap · cost · OR · routing",
    "                             ├─► validate.py   validation vs astram_event_data.csv (+ controlled check)",
    "                             └─► poi.py        OpenStreetMap context",
    "                                       │",
    "          build_all.py ───────────────►  outputs/  (csv · geojson · charts · metrics.json)",
    "                                       │",
    "                                  app.py (Streamlit) · routing.py (Mappls/OSRM) · mappls.py (OAuth)",
]):
    txt(mono, t, 10.5, INK, mono=True, space=2, first=(i == 0))
bullets(s, 0.95, 5.55, 11.7, 1.2, [
    "Boots from committed artifacts (no raw data needed) · 5-fold out-of-fold encoding prevents leakage · "
    "config-driven for any city · unit-tested (pytest).",
], size=12.5, space=5)
note(s, "One-liner: 'It's not a notebook — it's an engineered, reproducible system.' Mention leakage-free OOF "
        "encoding and unit tests if a technical judge asks. The app runs from pre-computed artifacts so the "
        "demo never needs the raw file.")

s = content("PRODUCT AT A GLANCE", "One screen — impact, KPIs, and ask-in-English",
            "The live command center: headline KPIs, natural-language query bar, and the congestion-impact map.")
img(s, 2.55, 1.95, 8.2, 4.55, HERO, caption="ParkSight command center — header KPIs + Ask ParkSight + Command Map",
    left=0.0)
note(s, "This is the 'first impression' slide. Point out the six KPI cards (298K violations, 2,534 hotspots / "
        "78 zones, 88.5% deploy coverage, congestion cost, 0.939 stability, 3.15× validation lift) and the "
        "Ask-in-English bar. Everything else in the deck drills into one of these.")

# ════════════════════════════ 02
divider("02", "DATA & METHODOLOGY", "Real partner data, a transparent impact score")
note(prs.slides[-1], "Section 2: the data we used and how we turn it into a defensible impact metric.")

s = content("DATA", "Only the provided datasets — no synthetic, no external training data")
card(s, 0.95, 1.9, 3.7, "298,445", "police parking violations", BLUE)
card(s, 4.85, 1.9, 3.7, "151 days", "Nov 2023 – Apr 2024", GREEN)
card(s, 8.75, 1.9, 3.6, "169 / 54", "junctions / stations", VIO)
bullets(s, 0.95, 3.75, 11.7, 2.9, [
    "Modelling (everything) — built ONLY on the provided Theme-1 file "
    "jan to may police violation_anonymized791b166.csv: geo-tagged violations (location, vehicle no./type, "
    "27 violation types, timestamps, officer & device IDs — all anonymised). 97% parking-related.",
    "Validation only — cross-checked against the provided event log astram_event_data.csv "
    "(congestion / accident / breakdown reports). Never used for training.",
    "Maps & routing — MapmyIndia (Mappls) + OpenStreetMap / OSRM fallback. Basemap & road-routing "
    "infrastructure, not a modelling dataset. Uber H3 (res 9, ≈150 m) is the spatial index.",
], size=14, space=11)
note(s, "Compliance point: ALL modelling uses only the given dataset. The second dataset is independent and used "
        "only to validate, which is exactly what makes the validation credible. Maps are partner infrastructure.")

s = content("METHODOLOGY", "Congestion-Impact Score (CIS)")
rect(s, 0.95, 1.85, 11.45, 0.9, NAVY, MSO_SHAPE.ROUNDED_RECTANGLE)
txt(box(s, 1.1, 2.02, 11.1, 0.6),
    "CIS  =  parking-volume  ×  mean severity  ×  chronicity  ×  peak-hour overlap     (scaled 0–100)",
    16, WHITE, True, PP_ALIGN.CENTER)
table(s, 0.95, 3.0, 11.45, [
    ["Component", "What it captures", "Source"],
    ["Volume", "How many violations occur in the cell", "counts per H3 cell"],
    ["Severity (0–1)", "How much the violation blocks traffic", "expert weights, cross-checked"],
    ["Chronicity", "Persistent vs one-off problem", "share of active days"],
    ["Peak overlap", "Concentration in rush hours", "fraction in peak windows"],
], col_w=[2.5, 6.35, 2.6], fs=12)
bullets(s, 0.95, 5.55, 11.7, 1.2, [
    "Leakage-free: 5-fold out-of-fold encoding + shrinkage (k=5). Honest: a controlled check shows CIS tracks "
    "congestion mainly via DENSITY — severity is an enforcement-triage layer, not a congestion predictor.",
], size=12.5, space=5)
note(s, "The CIS is the heart of the project: it turns 'a violation' into 'traffic impact'. Severity weights "
        "(main-road parking = 1.0, footpath = 0.4, document offences ≈ 0). Be honest about the controlled check "
        "— it shows intellectual rigour and pre-empts the 'isn't this just counting?' question.")

s = content("DETECTION", "Hotspots → enforcement zones (H3 + DBSCAN)")
bullets(s, 0.95, 1.85, 5.25, 4.5, [
    "Spatial index: Uber H3 resolution 9 (≈150 m hexagons).",
    "Aggregate each violation to its cell; compute CIS per cell.",
    "Cluster adjacent high-impact cells into ZONES with DBSCAN "
    "(haversine, eps≈350 m, min_samples=2 → genuine clusters).",
    "",
    f"{M['n_cells']:,} hotspot cells → {M['n_zones']} enforcement zones.",
    "~7% of cells hold 80% of total impact.",
    "Top zones are real chokepoints — Cottonpet, City Market, Chickpet — found with zero manual labels.",
], size=14, space=8)
img(s, 6.45, 1.85, 6.0, 4.55, CMAP, caption="Command Map — congestion-impact hexagons, zones & MapmyIndia route")
note(s, "Point to the map. The model independently rediscovers Cottonpet / City Market as #1 hotspots — "
        "places any Bengaluru traffic officer instantly recognises. That face-validity builds trust fast.")

s = content("VISUALISATION", "A 3-D congestion city — impact you can read at a glance",
            "Same impact surface, extruded: bar height = congestion-impact, so chokepoints literally stand out.")
img(s, 2.7, 1.95, 7.9, 4.45, CITY3D, caption="Command Map · 3D city (pydeck) on MapmyIndia tiles", accent=VIO)
note(s, "Quick 'wow' slide. The 3D extrusion makes the concentration obvious — a handful of towers dominate the "
        "skyline, which is the whole thesis (7% of cells = 80% of impact) shown visually.")

# ════════════════════════════ 03
divider("03", "MODELS & ACCURACY", "Forecasting · optimisation · validation")
note(prs.slides[-1], "Section 3 is for the technical judges: the models, the numbers, evaluated honestly.")

s = content("MODEL · FORECASTING", "Spatio-temporal risk model + uncertainty")
bullets(s, 0.95, 1.78, 5.3, 2.6, [
    "Target: expected violations per (cell × day × hour).",
    "Ensemble: HistGradientBoosting + LightGBM on log-intensity; 17 features "
    "(location, cyclic time, recency/trend, H3 neighbour density, severity).",
    "Uncertainty: quantile GBM → 80% prediction interval per cell.",
    "Validation: true TEMPORAL hold-out (train 16 wks → score 8 unseen wks).",
], size=13, space=8)
_gx, _gp = ho.get("holdout_r2_cellhour"), ho.get("persistence_r2_cellhour")
rect(s, 0.95, 4.55, 5.3, 1.92, MINT, MSO_SHAPE.ROUNDED_RECTANGLE, line=MINTL, lw=1)
rect(s, 0.95, 4.55, 0.1, 1.92, GREEN)
tf = box(s, 1.2, 4.66, 4.9, 1.72)
txt(tf, "What it adds over a 'same-as-last-month' baseline", 12, GREEN, True, first=True, space=4)
txt(tf, f"›  Beats persistence even at location×hour — R² {_gx} vs {_gp}.", 11, INK, space=3)
txt(tf, "›  Hour-level resolution + 80% confidence bands (persistence gives neither).", 11, INK, space=3)
txt(tf, f"›  Flags {emg.get('n_emerging', 278)} emerging cells a last-month copy can't see.", 11, INK, space=0)
img(s, 6.4, 1.85, 6.05, 4.55, FCAST, caption="Forecast — day×hour risk surface + accuracy metrics", accent=GREEN)
note(s, "Reframe the forecast as VALUE, not raw R². Yes, hotspots are persistent — but persistence gives one "
        "number per cell, no hours, and no uncertainty. Our model beats it even at location×hour, adds 80% "
        "confidence bands, and flags emerging cells a last-month copy structurally cannot see. That's the pitch.")

s = content("MODEL · OPTIMISATION", "Constrained patrol placement (Operations Research)")
bullets(s, 0.95, 1.85, 5.3, 4.5, [
    "Problem: place K patrols to cover the MOST impact — not just rank by score.",
    "Method: MAX-COVERAGE via submodular greedy (≈(1−1/e) guarantee).",
    "Each patrol covers its zone + zones within ≈1.2 km; impact counted once.",
    "",
    "Covers +~13 pp MORE unique impact than naive top-N — same budget — then sequenced into a "
    "road-following route via MapmyIndia.",
    "'Deploy & Simulate' slider recomputes coverage, route and ROI live.",
], size=13.5, space=7)
img(s, 6.45, 1.85, 6.0, 4.55, PRIO, caption="Priorities & Deploy — optimised vs naive coverage + outcome")
note(s, "This is the non-tautological intelligence: ranking by score stacks patrols on adjacent hotspots that "
        "overlap. Max-coverage spreads them to cover ~13 percentage points MORE unique impact with the same "
        "budget. That's a real OR contribution, not just analytics.")

s = content("ENFORCEMENT ANALYTICS", "The non-obvious intelligence: where you're NOT looking")
bullets(s, 0.95, 1.85, 5.3, 4.6, [
    f"Coverage-gap: {cov.get('n_under_enforced','13')} severe hotspots current patrols UNDER-serve "
    "(high impact, low effort).",
    "Enforcement-bias correction: estimate TRUE demand where patrols rarely go (+~30% hidden uplift).",
    f"Repeat offenders: {o['vio_from_repeat_%']}% of violations from repeat vehicles (worst {o['max_by_one_vehicle']}×).",
    f"Emerging hotspots: {emg.get('n_emerging','278')} rising spots flagged for early action.",
    "Ward-level unit allocation + congestion cost (₹, illustrative range).",
], size=13.5, space=8)
img(s, 6.45, 1.85, 6.0, 4.55, GAP, caption="Coverage & Context — under-enforced high-impact zones", accent=AMBER)
note(s, "Lead with COVERAGE-GAP — it's the killer, non-obvious insight: 'here are severe chokepoints your "
        "current patrols are MISSING.' That's intelligence beyond 'patrol where violations already are'.")

s = content("REPEAT OFFENDERS", "A small set of vehicles drives a large share of violations")
bullets(s, 0.95, 1.85, 5.3, 4.55, [
    f"{o['vio_from_repeat_%']}% of all violations come from repeat vehicles.",
    f"Worst single vehicle caught {o['max_by_one_vehicle']}× — chronic, targetable.",
    "Offender-concentration Pareto: the top few % of vehicles account for a disproportionate share.",
    "Deterrence check: do repeat offenders re-offend less after enforcement?",
    "Enables targeted notices / escalation — a lever beyond location-based patrols.",
], size=13.5, space=9)
img(s, 6.45, 1.85, 6.0, 4.55, OFF, caption="Repeat Offenders — concentration Pareto + chronic vehicles", accent=RED)
note(s, "Frames a second enforcement lever: not just WHERE, but WHO. A small set of habitual offenders is a "
        "concrete, actionable target for notices and escalation.")

s = content("LEARNING LOOP", "Post-enforcement tracking closes the loop")
bullets(s, 0.95, 1.85, 5.3, 4.55, [
    "Monthly improving vs worsening zones — measure whether action worked.",
    "City-wide violation trend over the full period.",
    "Animated hotspot evolution (monthly time-lapse).",
    "Emerging-hotspot watchlist for early intervention.",
    "This is the 'post-event learning system' the brief explicitly asks for.",
], size=13.5, space=9)
img(s, 6.45, 1.85, 6.0, 4.55, TREND, caption="Trends — improving / worsening / emerging + monthly trend", accent=GREEN)
note(s, "Tie directly to the problem statement's ask for a post-event learning system. Re-running monthly turns "
        "ParkSight from a snapshot into a feedback loop.")

s = content("WHY TRUST IT", "Validated against an INDEPENDENT dataset (astram_event_data.csv)")
card(s, 0.95, 1.75, 3.6, v["neighbourhood_pearson_res6"], "corr. with congestion density", GREEN, 1.2)
card(s, 4.75, 1.75, 3.6, f"{v['neighbourhood_events_top20pct_%']}%", "congestion in worst-20% areas", BLUE, 1.2)
card(s, 8.55, 1.75, 3.55, f"{v['neighbourhood_lift']}×", "more than random chance", RED, 1.2)
img(s, 0.95, 3.15, 5.6, 3.3, VAL, caption="Validation — congestion events by Impact-Score decile", accent=GREEN)
img(s, 6.85, 3.15, 5.55, 3.3, VALC, caption="Controlled check — partial correlation (honest)", accent=AMBER)
note(s, "THE differentiator: we validate the modelling dataset against a completely independent one "
        "(astram_event_data.csv). Worst-20% parking areas hold 63% of independently-reported congestion — 3.15× "
        "random. The controlled partial-correlation check is shown so we don't over-claim — associational, not causal.")

# ════════════════════════════ 04
divider("04", "PRODUCT & IMPACT", "A live command center, deployable today")
note(prs.slides[-1], "Section 4: show the product, the workflow, the impact, and how it deploys.")

s = content("THE PRODUCT", "A live command center — 9 interactive views")
feat = [("Command Map", "Impact hexagons · 3D city · MapmyIndia"), ("Ask ParkSight", "Natural-language queries"),
        ("Live Ops", "Highest-risk zones for THIS hour"), ("Priorities & Deploy", "Optimise · route · ROI · export"),
        ("Forecast", "Risk surface + confidence bands"), ("Repeat Offenders", "Offender Pareto · deterrence"),
        ("Trends", "Improving/worsening · time-lapse"), ("Coverage & Context", "Under-enforced · bias · POI")]
for i, (t, b) in enumerate(feat):
    cx = 0.95 + (i % 4) * 3.0; cy = 1.9 + (i // 4) * 2.35
    rect(s, cx, cy, 2.85, 2.1, LIGHT, MSO_SHAPE.ROUNDED_RECTANGLE, line=LINE, lw=1)
    rect(s, cx + 0.16, cy + 0.16, 0.42, 0.09, BLUE)
    tf = box(s, cx + 0.16, cy + 0.33, 2.55, 1.65)
    txt(tf, t, 13.5, INK, True, first=True); txt(tf, b, 11, MUTE, space=0)
txt(box(s, 0.95, 6.55, 11.7, 0.4), "Plus a one-click PDF Commander's Briefing and the independent Validation view.",
    12, MUTE)
note(s, "Don't read every card. Say: 'Nine views, one workflow — but three steal the show', then jump to the "
        "next slide / live demo.")

s = content("STANDOUT FEATURES", "What judges remember")
img(s, 0.95, 1.8, 5.75, 2.35, LIVE, caption="Live Ops — deploy THIS hour", accent=BLUE)
img(s, 6.85, 1.8, 5.55, 2.35, POI, caption="POI context — metro / market overlay", accent=VIO)
img(s, 0.95, 4.25, 5.75, 2.35, CITY3D, caption="3D congestion city (pydeck)", accent=GREEN)
bullets(s, 6.95, 4.3, 5.45, 2.3, [
    "› Ask-in-English deployment queries (offline)",
    "› 3D extruded impact map (pydeck)",
    "› Live-Ops 'right now' scheduling",
    "› MapmyIndia / Streets / Satellite basemaps",
    "› One-click PDF Commander's Briefing",
], size=12.5, space=6)
note(s, "If doing a live demo, THIS is the moment: type a question into Ask ParkSight, flip the 3D city, drag "
        "the deploy slider. 60 seconds of interaction beats five slides.")

s = content("USE CASE", "From data to deployment, in one shift")
for i, (n, t, b) in enumerate([
        ("1", "Morning brief", "DCP opens the dashboard / PDF briefing: today's top zones + ward allocation."),
        ("2", "Deploy", "Sets patrol count → optimiser picks zones, MapmyIndia routes them."),
        ("3", "Act now", "Live-Ops shows the highest-risk zones for THIS hour; officers reposition."),
        ("4", "Learn", "Next month, the tracker shows which zones improved — close the loop.")]):
    cy = 1.95 + i * 1.18
    rect(s, 0.95, cy, 0.7, 0.92, BLUE if i % 2 == 0 else NAVY, MSO_SHAPE.OVAL)
    txt(box(s, 0.95, cy + 0.22, 0.7, 0.5), n, 20, WHITE, True, PP_ALIGN.CENTER, first=True)
    tf = box(s, 1.95, cy + 0.04, 10.4, 1.0)
    txt(tf, t, 15, INK, True, first=True); txt(tf, b, 12.5, MUTE, space=0)
note(s, "Tell it as a story of a DCP's shift — morning brief → deploy → act this hour → learn next month. "
        "Makes the abstract analytics concrete and operational.")

s = content("CONTEXT", "Why hotspots happen — commercial & transit context (OpenStreetMap)")
bullets(s, 0.95, 1.85, 5.3, 4.5, [
    "Each top zone is matched to its nearest point of interest (metro, market, mall, bus station).",
    "72% of top hotspots sit within 400 m of a metro station, market or transit node.",
    "Confirms parking congestion concentrates around commercial & transit demand — exactly where supply is short.",
    "Turns 'a red hex' into an explainable, briefable cause an officer can act on.",
], size=14, space=9)
img(s, 6.45, 1.85, 6.0, 4.55, OSM, caption="Coverage & Context — nearest-POI matching for top zones", accent=VIO)
note(s, "Context makes the hotspots explainable and defensible. POIs come from OpenStreetMap (maps layer), not "
        "from any modelling dataset — keep that distinction clear if asked.")

s = content("IMPACT & ROI", "Focus beats blanket patrolling")
card(s, 0.95, 1.9, 3.7, "~7% → 80%", "cells cover 80% of impact", GREEN)
card(s, 4.85, 1.9, 3.7, f"₹{cost['annual_cost_cr_low']}–{cost['annual_cost_cr_high']} cr", "congestion cost / yr (est.)", RED)
card(s, 8.75, 1.9, 3.6, "+~13 pp", "optimiser vs naive coverage", VIO)
bullets(s, 0.95, 3.8, 11.7, 2.6, [
    "Enforce ~7% of the city → address 80% of parking-induced congestion impact.",
    "Max-coverage optimisation covers more impact per patrol than naive ranking.",
    "A what-if ROI estimates violations & rupees averted (adjustable effectiveness slider; illustrative).",
    "Deployable today — on the police's own data and MapmyIndia's own maps.",
], size=15, space=11)
note(s, "The memorable line: 'Enforce 7% of the city, fix 80% of the problem.' Frame the ₹ figure as an "
        "illustrative range, not a precise claim — judges respect the honesty.")

s = content("WHY PARKSIGHT STANDS OUT", "Against a typical hackathon entry")
table(s, 0.95, 1.9, 11.45, [
    ["Dimension", "Typical entry", "ParkSight"],
    ["Output", "A dashboard / heatmap", "A deployable patrol PLAN (zones, route, units)"],
    ["Impact", "Counts violations", "Quantifies traffic impact (severity-weighted)"],
    ["Trust", "Self-reported", "Validated vs an INDEPENDENT congestion dataset"],
    ["Method", "One model", "Forecast + OR optimisation + uncertainty"],
    ["Partners", "Generic OSS maps", "MapmyIndia maps + routing, BTP data"],
    ["Rigour", "Best-case numbers", "Temporal hold-out + honest limitations"],
], col_w=[2.3, 4.4, 4.75], fs=11.5, row_h=0.5)
note(s, "Use this as the closing argument before logistics. Each row is a reason to rank us above a "
        "pretty-dashboard entry. The 'Trust' and 'Partners' rows are the strongest.")

s = content("DEPLOYMENT & SCALABILITY", "Ready to pilot, easy to scale")
bullets(s, 0.95, 1.9, 5.5, 4.5, [
    "Plugs into existing workflow:",
    "›  Nightly violation export → ParkSight refresh → dashboard + field plans.",
    "›  Officers get today's route on a phone; commanders get ward allocation + PDF brief.",
    "›  Live public demo on Streamlit Cloud; secrets-managed Mappls keys.",
], size=14, space=8)
bullets(s, 6.65, 1.9, 5.7, 4.5, [
    "Scales with one config block:",
    "›  City centre, bounding box, H3 resolution, severity weights are config-driven.",
    "›  Same pipeline runs on any city's violation feed.",
    "›  Reproducible build (~100 s) + unit tests + artifact caching.",
], size=14, space=8)
note(s, "Answers the 'can this actually be used / scaled?' question. No new hardware — it consumes the data the "
        "police already collect. Multi-city is a config change, not a rewrite.")

s = content("TECH STACK & PARTNERS", "Built on the provided real-world infrastructure")
rect(s, 6.55, 1.9, 0.012, 4.4, LINE)
bullets(s, 0.95, 1.9, 5.4, 4.5, [
    "Partners:",
    "›  MapmyIndia (Mappls) — map tiles + India-grade road routing (OAuth, cached)",
    "›  Bengaluru Traffic Police (ASTraM) — violation & congestion datasets",
    "",
    "Models & methods:",
    "›  HistGradientBoosting + LightGBM ensemble (forecast)",
    "›  Quantile GBM (uncertainty) · DBSCAN (clustering)",
    "›  Submodular max-coverage (optimisation) · NNLS (calibration)",
], size=13.5, space=7)
bullets(s, 6.85, 1.9, 5.5, 4.5, [
    "Engineering:",
    "›  Python · pandas · NumPy · scikit-learn · LightGBM",
    "›  Uber H3 · Plotly · pydeck · Streamlit",
    "›  OSRM (routing fallback) · OpenStreetMap",
    "",
    "Rigour:",
    "›  5-fold out-of-fold encoding · temporal hold-out",
    "›  unit tests (pytest) · config-driven · reproducible build",
], size=13.5, space=7)
note(s, "Highlight that we USED the partner infrastructure (MapmyIndia + BTP data) — a scoring signal — and "
        "that routing/maps gracefully fall back to OSS if keys are absent.")

s = content("HOW IT WORKS", "Transparent method, assumptions, and limitations — in the app")
bullets(s, 0.95, 1.85, 5.3, 4.55, [
    "Every assumption is documented in the Method tab:",
    "›  Severity weights (expert-set, cross-checked)",
    "›  Cost model: ₹250 / vehicle-hour (conservative)",
    "›  Enforcement effectiveness: adjustable slider",
    "›  Spatial unit: H3 res 9 · Forecast: temporal hold-out",
    "Limitations stated openly (congestion proxy, selection bias, associational validation).",
], size=13, space=7)
img(s, 6.45, 1.85, 6.0, 4.55, METH, caption="Method tab — how ParkSight works + assumptions", accent=BLUE)
note(s, "Showing the in-app Method/limitations tab proves we're transparent by design, not just on a slide. "
        "Judges can audit every assumption themselves.")

s = content("LIMITATIONS & SCOPE", "What we're honest about")
bullets(s, 0.95, 1.9, 11.7, 4.4, [
    "›  We PROXY congestion (no live flow data provided); validated against an independent congestion log.",
    "›  Enforcement data = where patrols looked (selection bias) — surfaced and corrected.",
    "›  Temporal pattern = detection times (officer shifts) — framed as detection-risk, not raw demand.",
    "›  Validation is associational, not causal (controlled check shown).",
    "›  ₹ cost & enforcement effectiveness are illustrative ranges; forecast beats persistence by a small, honest edge.",
    "›  Privacy: all IDs anonymised. One city, ~5 months; pipeline is config-driven to generalise.",
], size=14.5, space=11)
txt(box(s, 0.95, 6.45, 11.7, 0.45),
    "Owning these is the point: a validated, deployable decision-support tool — not over-claimed.", 13, BLUE, True)
note(s, "Do NOT skip this slide. In a government room, the team that states its own limitations earns more "
        "trust than the team that over-claims. It signals maturity and makes every other number believable.")

# Close
s = prs.slides.add_slide(BLANK); _bg(s, NAVY); rect(s, 0, 0, 0.22, 7.5, BLUE)
txt(box(s, 0.9, 2.3, 11.5, 1.1), "ParkSight turns parking data into traffic relief.", 30, WHITE, True)
txt(box(s, 0.9, 3.6, 11.7, 1.4),
    "Detect · quantify · forecast · optimise · validate — complete, validated, deployable.  "
    "Enforce 7% of the city, fix 80% of the problem.", 18, RGBColor(0xC7, 0xD4, 0xEA))
txt(box(s, 0.9, 6.4, 11.7, 0.4),
    "Live demo + repository in the submission  ·  Powered by MapmyIndia & Bengaluru Traffic Police data",
    12, RGBColor(0x8A, 0xA0, 0xC0))
note(s, "Close on the one-liner and invite the live demo. Thank the partners (MapmyIndia + BTP). "
        "Have the public Streamlit link ready to open.")

for out in ("ParkSight_Deck.pptx", "ParkSight_Deck_v2.pptx", "ParkSight_Deck_final.pptx"):
    try:
        prs.save(str(ROOT / out)); print("saved", out, "with", len(prs.slides._sldIdLst), "slides"); break
    except PermissionError:
        print(out, "locked, trying next...")
